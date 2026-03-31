from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
CREAM = RGBColor(0xFD, 0xF9, 0xF0)
CREAM_DARK = RGBColor(0xF5, 0xED, 0xD8)
SAGE = RGBColor(0xA8, 0xC5, 0xA0)
SAGE_LIGHT = RGBColor(0xED, 0xF3, 0xE8)
HERBORA_GREEN = RGBColor(0x2E, 0x7D, 0x50)
HERBORA_DARK = RGBColor(0x1D, 0x5C, 0x38)
FOREST = RGBColor(0x2A, 0x5A, 0x3C)
TERRACOTTA = RGBColor(0xC4, 0x72, 0x4A)
TEXT_DARK = RGBColor(0x2D, 0x3B, 0x2E)
TEXT_MEDIUM = RGBColor(0x5A, 0x6B, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED_DARK = RGBColor(0xC0, 0x39, 0x2B)
GREEN_CHECK = RGBColor(0x27, 0xAE, 0x60)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=TEXT_DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=TEXT_DARK, icon="✓", icon_color=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(8)

        run_icon = p.add_run()
        run_icon.text = f"{icon}  "
        run_icon.font.size = Pt(font_size)
        run_icon.font.color.rgb = icon_color or HERBORA_GREEN
        run_icon.font.bold = True

        run_text = p.add_run()
        run_text.text = item
        run_text.font.size = Pt(font_size)
        run_text.font.color.rgb = color
        run_text.font.name = 'Calibri'
    return txBox


def add_comparison_row(slide, top, label, wp_text, herbora_text, wp_bad=True):
    # Label
    add_text(slide, Inches(0.8), top, Inches(2.5), Inches(0.5), label, 15, TEXT_DARK, True)
    # WP cell
    cell_wp = add_shape(slide, Inches(3.5), top, Inches(4.2), Inches(0.5), RGBColor(0xFD, 0xED, 0xED) if wp_bad else SAGE_LIGHT)
    cell_wp.text_frame.paragraphs[0].text = wp_text
    cell_wp.text_frame.paragraphs[0].font.size = Pt(13)
    cell_wp.text_frame.paragraphs[0].font.color.rgb = RED_DARK if wp_bad else TEXT_DARK
    cell_wp.text_frame.paragraphs[0].font.name = 'Calibri'
    cell_wp.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Herbora cell
    cell_h = add_shape(slide, Inches(7.9), top, Inches(4.6), Inches(0.5), RGBColor(0xE8, 0xF5, 0xE9))
    cell_h.text_frame.paragraphs[0].text = herbora_text
    cell_h.text_frame.paragraphs[0].font.size = Pt(13)
    cell_h.text_frame.paragraphs[0].font.color.rgb = HERBORA_GREEN
    cell_h.text_frame.paragraphs[0].font.name = 'Calibri'
    cell_h.text_frame.paragraphs[0].font.bold = True
    cell_h.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


# ============================
# SLIDE 1 - COVER
# ============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)

# Green accent bar top
add_rect(slide, 0, 0, W, Inches(0.08), HERBORA_GREEN)

# Left decorative bar
add_rect(slide, 0, 0, Inches(0.5), H, FOREST)

# Main content area
add_text(slide, Inches(1.5), Inches(1.2), Inches(10), Inches(0.6),
         "PRESENTAZIONE SITO WEB", 16, TERRACOTTA, True, PP_ALIGN.LEFT)

add_text(slide, Inches(1.5), Inches(1.9), Inches(10), Inches(1.5),
         "Herbora — Erboristeria", 52, HERBORA_DARK, True, PP_ALIGN.LEFT)

add_text(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(0.8),
         "Un sito su misura vs WordPress: perché la differenza conta", 26, TEXT_MEDIUM, False, PP_ALIGN.LEFT)

# Divider line
add_rect(slide, Inches(1.5), Inches(4.3), Inches(3), Inches(0.05), TERRACOTTA)

# Info
add_text(slide, Inches(1.5), Inches(4.8), Inches(5), Inches(0.4),
         "Sviluppato da Enrico Maria Caruso", 18, TEXT_MEDIUM, False)
add_text(slide, Inches(1.5), Inches(5.3), Inches(5), Inches(0.4),
         "enricomariacaruso.it", 16, SAGE, False)

# Right side - brand card
card = add_shape(slide, Inches(8.5), Inches(2), Inches(4), Inches(4), WHITE, SAGE)
add_text(slide, Inches(8.8), Inches(2.3), Inches(3.4), Inches(0.5),
         "herboralab.it", 24, HERBORA_GREEN, True, PP_ALIGN.CENTER)
add_text(slide, Inches(8.8), Inches(3), Inches(3.4), Inches(2.5),
         "Sito vetrina one-page\nHTML5 + CSS3 + JavaScript\nNessun CMS, nessun plugin\n100% codice proprietario\nOttimizzato per performance\ne conversione clienti", 16, TEXT_MEDIUM, False, PP_ALIGN.CENTER)

# Bottom bar
add_rect(slide, 0, Inches(7.3), W, Inches(0.2), HERBORA_GREEN)


# ============================
# SLIDE 2 - IL PROBLEMA DI WORDPRESS
# ============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, 0, 0, W, Inches(0.06), HERBORA_GREEN)
add_rect(slide, 0, 0, Inches(0.15), H, TERRACOTTA)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         "Il problema dei siti WordPress", 36, HERBORA_DARK, True)

add_text(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
         "Perché un sito WordPress generico non è la scelta migliore per un'attività come Herbora", 17, TEXT_MEDIUM)

# 4 problem cards
problems = [
    ("⚠️  Lentezza", "WordPress carica 30-80 file tra plugin,\ntemi e script. Tempi di caricamento\ndi 4-8 secondi = clienti persi."),
    ("🔓  Vulnerabilità", "Il 43% dei siti hackerati usa WordPress.\nPlugin non aggiornati = porta aperta\nper attacchi e malware."),
    ("📦  Pesantezza", "Temi generici con migliaia di funzioni\ninutili. Il 90% del codice caricato\nnon serve al tuo sito."),
    ("🔧  Manutenzione", "Aggiornamenti continui di WP, tema\ne plugin. Incompatibilità frequenti\nche rompono il sito."),
]

for i, (title, desc) in enumerate(problems):
    x = Inches(0.8 + (i % 2) * 6.2)
    y = Inches(1.9 + (i // 2) * 2.6)
    card = add_shape(slide, x, y, Inches(5.6), Inches(2.2), RGBColor(0xFD, 0xF0, 0xED), TERRACOTTA)
    add_text(slide, x + Inches(0.4), y + Inches(0.3), Inches(4.8), Inches(0.5),
             title, 22, RED_DARK, True)
    add_text(slide, x + Inches(0.4), y + Inches(0.9), Inches(4.8), Inches(1.2),
             desc, 15, TEXT_DARK)


# ============================
# SLIDE 3 - LA SOLUZIONE: SITO SU MISURA
# ============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_rect(slide, 0, 0, W, Inches(0.06), HERBORA_GREEN)
add_rect(slide, 0, 0, Inches(0.15), H, HERBORA_GREEN)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         "La soluzione: un sito su misura", 36, HERBORA_DARK, True)

add_text(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
         "Herbora merita un sito unico come i suoi prodotti: codice scritto a mano, zero sprechi, massime prestazioni", 17, TEXT_MEDIUM)

advantages = [
    ("⚡  Ultra-veloce", "Solo 3 file (HTML, CSS, JS) contro\n30-80 di WordPress. Caricamento\nin meno di 1 secondo."),
    ("🛡️  Sicuro al 100%", "Nessun database, nessun login admin,\nnessun plugin vulnerabile. Zero\nsuperficie di attacco."),
    ("🎨  Design unico", "Ogni pixel è pensato per Herbora.\nNessun template condiviso con\nmigliaia di altri siti."),
    ("💰  Zero costi ricorrenti", "Nessun hosting WordPress, nessun\nrinnovo tema/plugin premium.\nSolo hosting statico (gratis o pochi €/anno)."),
]

for i, (title, desc) in enumerate(advantages):
    x = Inches(0.8 + (i % 2) * 6.2)
    y = Inches(1.9 + (i // 2) * 2.6)
    card = add_shape(slide, x, y, Inches(5.6), Inches(2.2), WHITE, HERBORA_GREEN)
    add_text(slide, x + Inches(0.4), y + Inches(0.3), Inches(4.8), Inches(0.5),
             title, 22, HERBORA_GREEN, True)
    add_text(slide, x + Inches(0.4), y + Inches(0.9), Inches(4.8), Inches(1.2),
             desc, 15, TEXT_DARK)


# ============================
# SLIDE 4 - CONFRONTO DIRETTO
# ============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, 0, 0, W, Inches(0.06), HERBORA_GREEN)
add_rect(slide, 0, 0, Inches(0.15), H, HERBORA_GREEN)

add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
         "Confronto diretto: WordPress vs Herbora", 32, HERBORA_DARK, True)

# Header row
add_text(slide, Inches(0.8), Inches(1.1), Inches(2.5), Inches(0.5), "Caratteristica", 14, TEXT_MEDIUM, True)
hdr_wp = add_shape(slide, Inches(3.5), Inches(1.1), Inches(4.2), Inches(0.5), RED_DARK)
hdr_wp.text_frame.paragraphs[0].text = "❌  Sito WordPress generico"
hdr_wp.text_frame.paragraphs[0].font.size = Pt(14)
hdr_wp.text_frame.paragraphs[0].font.color.rgb = WHITE
hdr_wp.text_frame.paragraphs[0].font.bold = True
hdr_wp.text_frame.paragraphs[0].font.name = 'Calibri'
hdr_wp.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

hdr_h = add_shape(slide, Inches(7.9), Inches(1.1), Inches(4.6), Inches(0.5), HERBORA_GREEN)
hdr_h.text_frame.paragraphs[0].text = "✅  Sito Herbora su misura"
hdr_h.text_frame.paragraphs[0].font.size = Pt(14)
hdr_h.text_frame.paragraphs[0].font.color.rgb = WHITE
hdr_h.text_frame.paragraphs[0].font.bold = True
hdr_h.text_frame.paragraphs[0].font.name = 'Calibri'
hdr_h.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

rows = [
    ("Velocità caricamento", "3-8 secondi", "< 1 secondo"),
    ("File caricati", "30-80 file (plugin, temi, script)", "3 file (HTML + CSS + JS)"),
    ("Sicurezza", "Vulnerabile (database + login admin)", "Invulnerabile (sito statico, no DB)"),
    ("Design", "Template condiviso con migliaia di siti", "Design esclusivo, scritto a mano"),
    ("SEO / Google", "Codice sporco, lento = penalizzato", "Codice pulito, veloce = premiato"),
    ("Aggiornamenti", "WP + tema + plugin = rischio rotture", "Nessuno necessario, sempre stabile"),
    ("Costo hosting", "10-30€/mese (hosting WP)", "0-5€/mese (hosting statico / GitHub)"),
    ("Mobile", "Dipende dal tema (spesso mediocre)", "Responsive nativo, testato su ogni device"),
    ("GDPR / Cookie", "Plugin extra (spesso a pagamento)", "Integrato nativamente nel codice"),
    ("Proprietà codice", "Dipendi dal tema/plugin di terzi", "100% tuo, codice proprietario"),
]

for i, (label, wp, hb) in enumerate(rows):
    add_comparison_row(slide, Inches(1.7 + i * 0.55), label, wp, hb)


# ============================
# SLIDE 5 - PERFORMANCE / LIGHTHOUSE
# ============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_rect(slide, 0, 0, W, Inches(0.06), HERBORA_GREEN)
add_rect(slide, 0, 0, Inches(0.15), H, HERBORA_GREEN)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         "Performance: i numeri parlano chiaro", 36, HERBORA_DARK, True)

add_text(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
         "Google penalizza i siti lenti. Ogni secondo in più di caricamento = -7% conversioni (fonte: Google)", 17, TEXT_MEDIUM)

# Score circles
scores_wp = [("45", "Performance"), ("52", "Accessibilità"), ("60", "Best Practices"), ("48", "SEO")]
scores_hb = [("95+", "Performance"), ("92+", "Accessibilità"), ("95+", "Best Practices"), ("95+", "SEO")]

# WP section
add_text(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(0.5),
         "❌  WordPress medio (Lighthouse Score)", 18, RED_DARK, True)

for i, (score, label) in enumerate(scores_wp):
    x = Inches(1.0 + i * 1.45)
    y = Inches(2.5)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(1.1), Inches(1.1))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xED)
    circle.line.color.rgb = RED_DARK
    circle.line.width = Pt(3)
    tf = circle.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = score
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RED_DARK
    add_text(slide, x, y + Inches(1.2), Inches(1.1), Inches(0.4),
             label, 11, TEXT_MEDIUM, False, PP_ALIGN.CENTER)

# Herbora section
add_text(slide, Inches(7), Inches(1.9), Inches(5.5), Inches(0.5),
         "✅  Sito Herbora (Lighthouse Score)", 18, HERBORA_GREEN, True)

for i, (score, label) in enumerate(scores_hb):
    x = Inches(7.2 + i * 1.45)
    y = Inches(2.5)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(1.1), Inches(1.1))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
    circle.line.color.rgb = HERBORA_GREEN
    circle.line.width = Pt(3)
    tf = circle.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = score
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = HERBORA_GREEN
    add_text(slide, x, y + Inches(1.2), Inches(1.1), Inches(0.4),
             label, 11, TEXT_MEDIUM, False, PP_ALIGN.CENTER)

# Bottom impact section
add_rect(slide, Inches(0.8), Inches(4.5), Inches(11.7), Inches(0.05), SAGE)

impacts = [
    ("🚀  Tempo di caricamento", "WordPress: 4-8 sec\nHerbora: < 1 sec", "I visitatori abbandonano\ndopo 3 secondi di attesa"),
    ("📱  Mobile Score", "WordPress: ~50/100\nHerbora: 95+/100", "Il 65% del traffico\nviene da smartphone"),
    ("📈  Impatto su Google", "WordPress: penalizzato\nHerbora: premiato", "Siti veloci = posizioni\npiù alte su Google"),
]

for i, (title, data, impact) in enumerate(impacts):
    x = Inches(0.8 + i * 4.1)
    card = add_shape(slide, x, Inches(4.8), Inches(3.7), Inches(2.3), WHITE, SAGE)
    add_text(slide, x + Inches(0.3), Inches(4.9), Inches(3.2), Inches(0.5),
             title, 16, HERBORA_GREEN, True, PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.3), Inches(5.4), Inches(3.2), Inches(0.8),
             data, 14, TEXT_DARK, False, PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.3), Inches(6.2), Inches(3.2), Inches(0.6),
             impact, 12, TERRACOTTA, True, PP_ALIGN.CENTER)


# ============================
# SLIDE 6 - FUNZIONALITÀ SITO HERBORA
# ============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, 0, 0, W, Inches(0.06), HERBORA_GREEN)
add_rect(slide, 0, 0, Inches(0.15), H, HERBORA_GREEN)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         "Cosa include il sito Herbora", 36, HERBORA_DARK, True)

add_text(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
         "Ogni funzionalità è stata pensata per convertire visitatori in clienti", 17, TEXT_MEDIUM)

features = [
    ("🏠  9 Sezioni Complete",
     "Hero con blob organico\nStoria timeline interattiva\nFilosofia con layout editoriale\nServizi in griglia masonry\nCarosello prodotti\nPreparazioni su misura\nContatti con form e orari\nFooter completo"),
    ("🎨  Design Unico",
     "Palette botanica esclusiva\n(cream, salvia, terracotta)\nFont premium Google Fonts\nSVG botaniche disegnate a mano\nAnimazioni organiche clip-path\nDivisori onda tra sezioni\nNessun altro sito è uguale"),
    ("📱  Mobile First",
     "Menu mobile con overlay botanico\nLayout responsive su ogni device\nTouch-friendly (swipe, tap)\nImmagini ottimizzate\nBottoni grandi e accessibili\nTestato su Chrome DevTools\nsu tutti i breakpoint"),
    ("⚙️  Funzionalità Avanzate",
     'Indicatore "Aperto/Chiuso" live\nWhatsApp click-to-chat\nForm contatto con floating labels\nGoogle Maps (GDPR compliant)\nCookie consent con modal\nSchema.org per Google\nSEO meta tags + Open Graph'),
]

for i, (title, items) in enumerate(features):
    x = Inches(0.6 + (i % 4) * 3.15)
    y = Inches(1.7)
    card = add_shape(slide, x, y, Inches(2.9), Inches(5.3), CREAM, SAGE)
    add_text(slide, x + Inches(0.2), y + Inches(0.2), Inches(2.5), Inches(0.6),
             title, 16, HERBORA_GREEN, True, PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.2), y + Inches(0.9), Inches(2.5), Inches(4.2),
             items, 13, TEXT_DARK, False, PP_ALIGN.LEFT)


# ============================
# SLIDE 7 - COSTI A CONFRONTO
# ============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_rect(slide, 0, 0, W, Inches(0.06), HERBORA_GREEN)
add_rect(slide, 0, 0, Inches(0.15), H, HERBORA_GREEN)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         "Costi: WordPress vs Sito su misura", 36, HERBORA_DARK, True)

add_text(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
         "Oltre al costo iniziale, WordPress ha costi nascosti che si accumulano anno dopo anno", 17, TEXT_MEDIUM)

# WordPress costs
wp_card = add_shape(slide, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.2), WHITE, RED_DARK)
add_text(slide, Inches(1.1), Inches(1.9), Inches(5), Inches(0.5),
         "❌  WordPress — Costi annuali", 20, RED_DARK, True, PP_ALIGN.CENTER)

wp_costs = [
    "Hosting WordPress:  120-360€/anno",
    "Tema premium:  50-80€/anno",
    "Plugin premium (SEO, cache, sicurezza):  100-300€/anno",
    "Certificato SSL:  0-50€/anno",
    "Manutenzione/aggiornamenti:  200-500€/anno",
    "Backup e sicurezza:  50-150€/anno",
    "Fixing problemi/incompatibilità:  imprevisto",
]
add_bullet_list(slide, Inches(1.2), Inches(2.6), Inches(5), Inches(3.5),
                wp_costs, 13, TEXT_DARK, "•", RED_DARK)

add_text(slide, Inches(1.1), Inches(6.1), Inches(5), Inches(0.5),
         "Totale:  520 - 1.440€/anno", 20, RED_DARK, True, PP_ALIGN.CENTER)

# Herbora costs
hb_card = add_shape(slide, Inches(6.9), Inches(1.8), Inches(5.6), Inches(5.2), WHITE, HERBORA_GREEN)
add_text(slide, Inches(7.2), Inches(1.9), Inches(5), Inches(0.5),
         "✅  Sito su misura — Costi annuali", 20, HERBORA_GREEN, True, PP_ALIGN.CENTER)

hb_costs = [
    "Hosting (GitHub Pages):  GRATIS",
    "Tema/plugin:  NON SERVONO",
    "Certificato SSL:  INCLUSO (gratis)",
    "Manutenzione:  MINIMA (nessun aggiornamento)",
    "Sicurezza:  NATIVA (nessun database)",
    "Dominio personalizzato:  10-15€/anno",
    "Modifiche contenuti:  su richiesta",
]
add_bullet_list(slide, Inches(7.3), Inches(2.6), Inches(5), Inches(3.5),
                hb_costs, 13, TEXT_DARK, "✓", HERBORA_GREEN)

add_text(slide, Inches(7.2), Inches(6.1), Inches(5), Inches(0.5),
         "Totale:  10 - 15€/anno", 20, HERBORA_GREEN, True, PP_ALIGN.CENTER)


# ============================
# SLIDE 8 - GDPR E CONFORMITÀ
# ============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, 0, 0, W, Inches(0.06), HERBORA_GREEN)
add_rect(slide, 0, 0, Inches(0.15), H, HERBORA_GREEN)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         "GDPR, Privacy e Conformità", 36, HERBORA_DARK, True)

add_text(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
         "Il sito Herbora è conforme al GDPR fin dalla progettazione (Privacy by Design)", 17, TEXT_MEDIUM)

gdpr_features = [
    ("🍪  Cookie Consent", "Banner GDPR con 3 opzioni:\nAccetta tutti / Solo necessari / Personalizza\nModal dettagliato per categoria\nPreferenze salvate per 6 mesi\nNessun cookie caricato senza consenso"),
    ("🗺️  Google Maps condizionale", "La mappa si carica SOLO se l'utente\naccetta i cookie di terze parti.\nIn caso contrario, viene mostrato un\nplaceholder con link diretto a Google Maps."),
    ("📋  Pagine legali complete", "Privacy Policy completa\n(titolare, finalità, diritti, terze parti)\nCookie Policy dettagliata\ncon tabella cookie utilizzati\nLink ai browser per gestione cookie"),
    ("🔒  Sicurezza nativa", "Nessun database = nessun data breach\nNessun form di login admin\nDati form inviati via FormSubmit\n(nessun dato salvato sul server)\nHTTPS forzato su GitHub Pages"),
]

for i, (title, desc) in enumerate(gdpr_features):
    x = Inches(0.6 + (i % 2) * 6.2)
    y = Inches(1.8 + (i // 2) * 2.7)
    card = add_shape(slide, x, y, Inches(5.8), Inches(2.4), CREAM, SAGE)
    add_text(slide, x + Inches(0.3), y + Inches(0.2), Inches(5.2), Inches(0.5),
             title, 19, HERBORA_GREEN, True)
    add_text(slide, x + Inches(0.3), y + Inches(0.8), Inches(5.2), Inches(1.5),
             desc, 14, TEXT_DARK)


# ============================
# SLIDE 9 - RIEPILOGO VANTAGGI
# ============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, FOREST)
add_rect(slide, 0, 0, W, Inches(0.06), TERRACOTTA)

add_text(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
         "Perché scegliere il sito su misura", 36, WHITE, True)

add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
         "Tutti i vantaggi in sintesi", 18, SAGE, False)

benefits = [
    "Velocità superiore: caricamento in meno di 1 secondo (vs 4-8 sec WordPress)",
    "Sicurezza totale: nessun database, nessun plugin, nessuna vulnerabilità",
    "Design esclusivo: un sito che rappresenta davvero l'identità di Herbora",
    "SEO ottimizzato: codice pulito, Schema.org, meta tags = più visibilità su Google",
    "Zero manutenzione: nessun aggiornamento, nessun rischio di rotture",
    "Costi minimi: 10-15€/anno vs 500-1.400€/anno di WordPress",
    "GDPR nativo: cookie consent, privacy policy, tutto integrato e a norma",
    "Mobile perfetto: responsive nativo, testato su ogni dispositivo",
    "WhatsApp integrato: i clienti ti contattano con un tap",
    "Proprietà totale: il codice è tuo, nessuna dipendenza da terzi",
]

add_bullet_list(slide, Inches(1.2), Inches(1.9), Inches(11), Inches(5),
                benefits, 17, WHITE, "✓", SAGE)


# ============================
# SLIDE 10 - CONTATTI / CHIUSURA
# ============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_rect(slide, 0, 0, W, Inches(0.06), HERBORA_GREEN)

# Main content centered
add_text(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(1),
         "Il tuo sito è online e pronto.", 42, HERBORA_DARK, True, PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(2.7), Inches(11.3), Inches(0.6),
         "enrix84.github.io/herbora", 28, TERRACOTTA, True, PP_ALIGN.CENTER)

add_rect(slide, Inches(5.5), Inches(3.6), Inches(2.3), Inches(0.04), SAGE)

add_text(slide, Inches(1), Inches(4), Inches(11.3), Inches(1),
         "Sviluppato da", 18, TEXT_MEDIUM, False, PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(4.5), Inches(11.3), Inches(0.8),
         "Enrico Maria Caruso", 32, HERBORA_DARK, True, PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(5.3), Inches(11.3), Inches(0.5),
         "Full Stack Developer — enricomariacaruso.it", 18, SAGE, False, PP_ALIGN.CENTER)

# Bottom bar
add_rect(slide, 0, Inches(7.3), W, Inches(0.2), HERBORA_GREEN)


# SAVE
prs.save('C:/workspace/herbora/Herbora_Presentazione.pptx')
print("Presentazione creata con successo!")
